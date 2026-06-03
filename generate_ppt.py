import os
from pptx import Presentation
from pptx.util import Inches, Pt
from bs4 import BeautifulSoup
import glob

def create_presentation(directory, output_file):
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Atharve Portfolio"
    subtitle.text = "Website Template Overview"
    
    html_files = glob.glob(os.path.join(directory, "*.html"))
    html_files = [f for f in html_files if "node_modules" not in f]
    
    # Sort files to put index.html first
    html_files.sort(key=lambda x: 0 if os.path.basename(x) == "index.html" else 1)
    
    for file_path in html_files:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                
            page_title = soup.title.string if soup.title else os.path.basename(file_path)
            if not page_title:
                page_title = os.path.basename(file_path)
                
            # Create a slide for each page
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = f"{page_title.strip()}"
            
            tf = body_shape.text_frame
            tf.text = f"File: {os.path.basename(file_path)}"
            
            # Extract headings
            headings = soup.find_all(['h1', 'h2', 'h3'])
            if headings:
                p = tf.add_paragraph()
                p.text = "Key Sections:"
                p.level = 0
                for i, h in enumerate(headings[:8]):  # limit to first 8 headings
                    text = h.get_text().strip()
                    if text:
                        p = tf.add_paragraph()
                        p.text = text
                        p.level = 1
            else:
                p = tf.add_paragraph()
                p.text = "No major headings found."
                p.level = 0
                
        except Exception as e:
            print(f"Error processing {file_path}: {e}")

    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    create_presentation(".", "Atharve_Portfolio_Overview.pptx")
